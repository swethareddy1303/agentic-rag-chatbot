# utils/file_parser.py

import fitz  # PyMuPDF
import pandas as pd
import docx
import pptx
import os

def parse_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    elif ext in [".txt", ".md"]:
        return parse_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def parse_pdf(path):
    text = ""
    doc = fitz.open(path)
    for page in doc:
        text += page.get_text()
    return text

def parse_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def parse_pptx(path):
    prs = pptx.Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_csv(path):
    df = pd.read_csv(path)
    return df.to_string(index=False)

def parse_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()
